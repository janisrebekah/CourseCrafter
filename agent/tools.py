
from docx import Document
import os
import json
from typing import TypedDict, List
from io import BytesIO
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
import logging
logging.basicConfig(level=logging.INFO)

#final word
def generate_quiz_docx(course_structure: dict, output_filename: str = "quiz.docx") -> str:
    doc = Document()
    doc.add_heading(course_structure.get("courseTitle", "Course Quiz"), level=0)

    for i, module in enumerate(course_structure.get("modules", []), 1):
        doc.add_heading(f"Module {i}: {module.get('moduleTitle')}", level=1)
        for q in range(1, 4):
            doc.add_paragraph(f"Q{q}. Write a short answer on: {module.get('topics')[q % len(module.get('topics'))]}")

    output_dir = os.path.join(os.path.dirname(__file__), '..', 'output')
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, output_filename)
    doc.save(output_path)
    return f"Quiz Word document created: {output_path}"

def create_dynamic_presentation(ppt_spec: dict, 
                                output_filename: str = "dynamic_presentation.pptx") -> str:
    from pathlib import Path
    prs = Presentation()

    # Create a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = ppt_spec.get("title", "Untitled Presentation")
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Generated PPT for course module"

    for spec in ppt_spec.get("slides", []):
        layout_type = spec.get("layout", "default")

        if layout_type == "title_only":
            slide_layout = prs.slide_layouts[5]
        else:
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = spec.get("title", "Untitled Slide")

        if "content" in spec:
            placeholder = slide.placeholders[1]
            text_frame = placeholder.text_frame
            text_frame.clear()
            text_frame.add_paragraph().text = spec["content"]

        if "bullets" in spec:
            placeholder = slide.placeholders[1]
            text_frame = placeholder.text_frame
            text_frame.clear()
            for bullet in spec["bullets"]:
                text_frame.add_paragraph().text = bullet

        if "image" in spec:
            try:
                image_source = spec["image"]
                if image_source.startswith("http"):
                    response = requests.get(image_source)
                    image_data = BytesIO(response.content)
                else:
                    image_data = image_source
                slide.shapes.add_picture(image_data, Inches(5), Inches(1.5), width=Inches(3))
            except Exception as e:
                logging.warning(f"Image error: {e}")

    # ✅ SAVE FILE HERE
    output_path = Path(__file__).resolve().parent.parent / "output" / output_filename
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    logging.info(f"Presentation created at: {output_path}")
    return f"✅ PPT saved: {output_path}"
